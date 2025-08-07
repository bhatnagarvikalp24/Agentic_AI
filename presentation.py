"""
Presentation utilities for the Agentic AI Assistant.
Contains PowerPoint generation functionality and related utilities.
"""

import pandas as pd
import os
import re
from io import BytesIO
from typing import Dict, Any, Union, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor


def generate_ppt(entry: Dict[str, Any]) -> BytesIO:
    """
    Generate a PowerPoint presentation from chat entry data.
    
    Args:
        entry: Dictionary containing chat entry data with route, results, etc.
        
    Returns:
        BytesIO object containing the PowerPoint file
    """
    prs = Presentation()
    layout = prs.slide_layouts[5]  # title + content

    # 🧠 Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Agentic AI Report"
    slide.placeholders[1].text = f"Prompt: {entry['prompt']}"

    route = entry.get("route")

    # 🧾 SQL Query Slide (if applicable)
    if route in ["sql", "document", "comp"] and entry.get("sql_query"):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "SQL Query"
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = entry["sql_query"]
        p.font.size = Pt(14)

    # 📊 SQL Result Table (if applicable)
    result = entry.get("result")
    if isinstance(result, list):
        result = pd.DataFrame(result)

    if route in ["sql", "document", "comp"] and isinstance(result, pd.DataFrame) and not result.empty:
        df = pd.DataFrame(entry["result"]) if isinstance(entry["result"], list) else entry["result"]
        if isinstance(df, pd.DataFrame):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = "SQL Results"
            rows = min(6, len(df) + 1)
            cols = len(df.columns)
            table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(8.5), Inches(3)).table
            for i, col in enumerate(df.columns):
                table.cell(0, i).text = str(col)
            for i, row in df.head(5).iterrows():
                for j, val in enumerate(row):
                    table.cell(i + 1, j).text = str(val)

    # 🆚 Comparison Summary
    if route == "comp" and entry.get("comparison_summary"):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Comparison Summary"
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        for para in entry["comparison_summary"].split("\n"):
            if para.strip():
                p = tf.add_paragraph()
                p.text = para.strip()
                p.font.size = Pt(14)
                p.space_after = Pt(4)

    # 🧠 General Summary (Search + Comp)
    if route in ["search", "comp"] and entry.get("general_summary"):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "General Summary"
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        for para in entry["general_summary"].split("\n"):
            if para.strip():
                p = tf.add_paragraph()
                p.text = para.strip()
                p.font.size = Pt(14)
                p.space_after = Pt(4)

    # 🔗 Top Web Links (Search + Comp)
    if route in ["search", "comp"] and entry.get("web_links"):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Top Web Links"
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True

        for i, (link_md, summary) in enumerate(entry["web_links"], 1):
            # Match Markdown-style link: [Title](https://link)
            match = re.match(r"\[(.*?)\]\((.*?)\)", link_md)
            if match:
                title, url = match.groups()
            else:
                title, url = f"Link {i}", link_md  # fallback

            # Add hyperlink paragraph
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"{i}. {title}"
            run.font.size = Pt(13)
            run.hyperlink.address = url
            p.space_after = Pt(2)

            # Add summary (not a hyperlink)
            summary_p = tf.add_paragraph()
            summary_p.text = f"    ↳ {summary[:180]}..."
            summary_p.font.size = Pt(12)
            summary_p.space_after = Pt(6)

    # 📘 FAISS Knowledge Base Slides
    if route == "faissdb":
        _add_faiss_slides(prs, entry, layout)

    # Finalize PPT in memory
    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes


def _add_faiss_slides(prs: Presentation, entry: Dict[str, Any], layout) -> None:
    """
    Add FAISS knowledge base slides to the presentation.
    
    Args:
        prs: Presentation object
        entry: Chat entry data
        layout: Slide layout to use
    """
    # 🧠 faiss Summary Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "FAISS Summary"

    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True

    summary_text = entry.get("faiss_summary", "No summary available.")
    for para in summary_text.split("\n"):
        if para.strip():
            p = tf.add_paragraph()
            p.text = para.strip()
            p.font.size = Pt(14)
            p.space_after = Pt(4)

    # 📄 Source Slides (with clickable file name in title if available)
    for i, (docname, snippet, path) in enumerate(entry.get("faiss_sources", []), 1):
        filename = os.path.basename(path) if path else docname
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Source {i}: {filename}"

        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(5.0))
        tf = box.text_frame
        tf.word_wrap = True

        for para in snippet.split("\n"):
            if para.strip():
                p = tf.add_paragraph()
                p.text = para.strip()
                p.font.size = Pt(12)
                p.space_after = Pt(3)

    # 🖼️ Image Slide (only from the most similar document)
    faiss_images = entry.get("faiss_images", [])
    faiss_sources = entry.get("faiss_sources", [])
    if faiss_images and faiss_sources:
        top_docname = faiss_sources[0][0]  # first doc's name
        top_doc_images = [img for img in faiss_images if img.get("original_doc") == top_docname]

        if top_doc_images:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
            slide.shapes.title.text = f"Images from {top_docname}"

            # Adjust layout
            left = Inches(0.8)
            top = Inches(2.5)  # ⬅️ Increased top margin
            image_width = Inches(6)  # ⬅️ Increased width
            spacing = Inches(0.6)

            for idx, img_meta in enumerate(top_doc_images):
                img_path = img_meta.get("extracted_image_path")
                if img_path and os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, left, top, width=image_width)
                    top += Inches(3.2)  # ⬅️ Increased vertical spacing

                    if top > Inches(6.5):  # wrap to next column if needed
                        top = Inches(2.0)
                        left += image_width + spacing


def create_simple_ppt(title: str, content: str) -> BytesIO:
    """
    Create a simple PowerPoint with title and content.
    
    Args:
        title: Presentation title
        content: Content text
        
    Returns:
        BytesIO object containing the PowerPoint file
    """
    prs = Presentation()
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    
    # Content slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(14)
    
    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes


def add_table_to_slide(slide, data: List[List[str]], title: str = "") -> None:
    """
    Add a table to a slide.
    
    Args:
        slide: Slide object to add table to
        data: 2D list of data for the table
        title: Optional title for the slide
    """
    if title:
        slide.shapes.title.text = title
    
    if not data:
        return
    
    rows = len(data)
    cols = len(data[0]) if data else 0
    
    if rows > 0 and cols > 0:
        table = slide.shapes.add_table(
            rows, cols, 
            Inches(0.5), Inches(1.2), 
            Inches(8.5), Inches(3)
        ).table
        
        for i, row in enumerate(data):
            for j, cell_value in enumerate(row):
                if i < len(table.rows) and j < len(table.columns):
                    table.cell(i, j).text = str(cell_value)


def add_text_to_slide(slide, text: str, title: str = "") -> None:
    """
    Add text content to a slide.
    
    Args:
        slide: Slide object to add text to
        text: Text content to add
        title: Optional title for the slide
    """
    if title:
        slide.shapes.title.text = title
    
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    
    for para in text.split("\n"):
        if para.strip():
            p = tf.add_paragraph()
            p.text = para.strip()
            p.font.size = Pt(14)
            p.space_after = Pt(4)


# Constants for presentation styling
SLIDE_TITLES = {
    "sql": "SQL Query Results",
    "search": "Web Search Results", 
    "document": "Document Update",
    "comp": "Comparison Analysis",
    "faissdb": "Knowledge Base Results"
}

DEFAULT_FONT_SIZE = 14
DEFAULT_SPACING = 4 